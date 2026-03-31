"""Fast parsers for when Docling is too slow or unnecessary.

Each parser extracts text + basic metadata and chunks by character count.
Used in 'quality' mode (for large files / XLSX) and 'fast' mode (always).
"""

from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path

from stripes_rag.chunker import ChunkResult
from stripes_rag.config import settings


def _chunk_text(
    text: str,
    headings: str | None = None,
    page_numbers: str | None = None,
    start_index: int = 0,
) -> list[ChunkResult]:
    """Split text into chunks by character count with smart split points."""
    max_chars = settings.chunk_max_tokens * 4
    if not text.strip():
        return []

    if len(text) <= max_chars:
        return [ChunkResult(
            text=text.strip(),
            headings=headings,
            page_numbers=page_numbers,
            chunk_index=start_index,
        )]

    results: list[ChunkResult] = []
    remaining = text
    idx = start_index

    while remaining:
        if len(remaining) <= max_chars:
            if remaining.strip():
                results.append(ChunkResult(
                    text=remaining.strip(),
                    headings=headings,
                    page_numbers=page_numbers,
                    chunk_index=idx,
                ))
            break

        # Find best split point within max_chars
        chunk = remaining[:max_chars]
        split_at = -1
        for sep in ["\n\n", "\n", ". ", " "]:
            pos = chunk.rfind(sep)
            if pos > max_chars // 4:  # don't split too early
                split_at = pos + len(sep)
                break
        if split_at == -1:
            split_at = max_chars  # hard cut

        piece = remaining[:split_at].strip()
        if piece:
            results.append(ChunkResult(
                text=piece,
                headings=headings,
                page_numbers=page_numbers,
                chunk_index=idx,
            ))
            idx += 1
        remaining = remaining[split_at:]

    return results


# ---------------------------------------------------------------------------
# Per-format parsers
# ---------------------------------------------------------------------------

def _parse_xlsx(file_path: Path, file_hash: str) -> list[ChunkResult]:
    import pandas as pd

    sheets = pd.read_excel(file_path, sheet_name=None, dtype=str)
    results: list[ChunkResult] = []
    idx = 0

    for sheet_name, df in sheets.items():
        df = df.dropna(how="all").dropna(axis=1, how="all")
        if df.empty:
            continue
        md = df.to_markdown(index=False)
        heading = str(sheet_name)
        chunks = _chunk_text(md, headings=heading, start_index=idx)
        results.extend(chunks)
        idx += len(chunks)

    return results


def _parse_pdf(file_path: Path, file_hash: str) -> list[ChunkResult]:
    import fitz

    doc = fitz.open(file_path)
    results: list[ChunkResult] = []
    idx = 0

    for page_num in range(len(doc)):
        page = doc[page_num]
        text = page.get_text()
        if not text.strip():
            continue
        page_str = str(page_num + 1)
        chunks = _chunk_text(text, page_numbers=page_str, start_index=idx)
        results.extend(chunks)
        idx += len(chunks)

    doc.close()
    return results


def _parse_docx(file_path: Path, file_hash: str) -> list[ChunkResult]:
    from docx import Document

    doc = Document(str(file_path))
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    text = "\n\n".join(paragraphs)
    return _chunk_text(text)


def _parse_pptx(file_path: Path, file_hash: str) -> list[ChunkResult]:
    from pptx import Presentation

    prs = Presentation(str(file_path))
    results: list[ChunkResult] = []
    idx = 0

    for slide_num, slide in enumerate(prs.slides, 1):
        parts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    if para.text.strip():
                        parts.append(para.text)
        if parts:
            text = "\n".join(parts)
            chunks = _chunk_text(
                text,
                headings=f"Slide {slide_num}",
                page_numbers=str(slide_num),
                start_index=idx,
            )
            results.extend(chunks)
            idx += len(chunks)

    return results


def _parse_html(file_path: Path, file_hash: str) -> list[ChunkResult]:
    from html.parser import HTMLParser

    class _TextExtractor(HTMLParser):
        def __init__(self):
            super().__init__()
            self.parts: list[str] = []
            self._skip = False

        def handle_starttag(self, tag, attrs):
            self._skip = tag in ("script", "style")

        def handle_endtag(self, tag):
            if tag in ("script", "style"):
                self._skip = False

        def handle_data(self, data):
            if not self._skip and data.strip():
                self.parts.append(data.strip())

    raw = file_path.read_text(errors="replace")
    parser = _TextExtractor()
    parser.feed(raw)
    text = "\n".join(parser.parts)
    return _chunk_text(text)


def _parse_md(file_path: Path, file_hash: str) -> list[ChunkResult]:
    text = file_path.read_text(errors="replace")
    return _chunk_text(text)


_PARSERS = {
    ".xlsx": _parse_xlsx,
    ".pdf": _parse_pdf,
    ".docx": _parse_docx,
    ".pptx": _parse_pptx,
    ".html": _parse_html,
    ".md": _parse_md,
}


def fast_parse_and_chunk(file_path: Path, file_hash: str) -> list[ChunkResult]:
    """Parse and chunk a file using fast parsers (no Docling)."""
    ext = file_path.suffix.lower()
    parser = _PARSERS.get(ext)
    if parser is None:
        raise ValueError(f"No fast parser for {ext}")
    return parser(file_path, file_hash)
